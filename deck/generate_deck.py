#!/usr/bin/env python3
"""BambooHR case study deck — problem-led story + client / BambooHR why + process."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
SHOTS = ROOT / "screenshots"
LOGO = SHOTS / "bamboohr-logo.png"
SPREADSHEET = SHOTS / "00-current-spreadsheet.png"

GREEN = RGBColor(0x73, 0xC4, 0x1D)
GREEN_DK = RGBColor(0x4F, 0x8F, 0x12)
INK = RGBColor(0x1B, 0x1F, 0x17)
MUTED = RGBColor(0x6A, 0x72, 0x64)
LINE = RGBColor(0xE6, 0xEA, 0xE1)
SOFT = RGBColor(0xF3, 0xF8, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xFA, 0xFB, 0xF8)

TOTAL = 14
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def font(run, size=18, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"


def text(shape, content, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    font(r, size, bold, color)
    return tf


def para(tf, content, size=14, bold=False, color=INK, before=8, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    r = p.add_run()
    r.text = content
    font(r, size, bold, color)
    return p


def slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, l, t, w, h, fill=WHITE, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    return sh


def bar(s):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    sh.fill.solid()
    sh.fill.fore_color.rgb = GREEN
    sh.line.fill.background()


def logo(s, left=Inches(0.75), top=Inches(0.28), width=Inches(1.55)):
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), left, top, width=width)


def foot(s, n):
    t = s.shapes.add_textbox(Inches(0.75), Inches(7.15), Inches(10), Inches(0.25))
    text(t, "BambooHR interview case study", size=11, color=MUTED)
    p = s.shapes.add_textbox(Inches(12.2), Inches(7.15), Inches(0.8), Inches(0.25))
    text(p, f"{n}/{TOTAL}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def notes(s, body):
    s.notes_slide.notes_text_frame.text = body


def head(s, kicker, title, sub=None):
    bar(s)
    logo(s)
    k = s.shapes.add_textbox(Inches(2.55), Inches(0.25), Inches(10), Inches(0.25))
    text(k, kicker.upper(), size=11, bold=True, color=GREEN_DK)
    t = s.shapes.add_textbox(Inches(2.55), Inches(0.5), Inches(10.2), Inches(0.55))
    text(t, title, size=26, bold=True, color=INK)
    if sub:
        u = s.shapes.add_textbox(Inches(2.55), Inches(1.05), Inches(10.2), Inches(0.35))
        text(u, sub, size=15, color=MUTED)


# ---------- 1 Title ----------
s = slide()
wash = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
wash.fill.solid()
wash.fill.fore_color.rgb = WASH
wash.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(7.5))
accent.fill.solid()
accent.fill.fore_color.rgb = GREEN
accent.line.fill.background()
if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.9), Inches(1.5), width=Inches(2.4))
t = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11), Inches(1.6))
tf = text(t, "Headcount Control Tower", size=42, bold=True, color=INK)
para(tf, "Product discovery case study", size=20, color=MUTED, before=14)
m = s.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11), Inches(1.4))
tf = text(m, "The problem: leaders can’t trust or explain planned vs filled headcount.", size=18, color=INK)
para(tf, "The bet: make that story clear in 30 seconds — for Finance, HR, and BambooHR.", size=18, color=INK, before=10)
para(tf, "AI-augmented discovery  ·  Interactive prototype  ·  Partner review with Head of Finance", size=14, color=MUTED, before=12)
notes(s, "Open on the problem. Prototype is the proof; this deck is why it matters.")


# ---------- 2 Problem we solve ----------
s = slide()
head(
    s,
    "The problem we solve",
    "Reconciliation is critical — but nobody can follow the story.",
    "Not a missing-data problem. A comprehension and alignment problem.",
)
pains = [
    (
        "The number isn’t trusted",
        "Approved plan, working plan, and filled headcount live in different tabs and systems. Leaders argue whose number is right.",
    ),
    (
        "The gap isn’t explained",
        "Marketing can be −2 vs plan with no named reason — not backfilling, role closed, or a pivot to another team.",
    ),
    (
        "The date isn’t clear",
        "Is this as of today or as of month-end close? Without a trusted as-of, every recon rebuilds trust from scratch.",
    ),
]
for i, (h, d) in enumerate(pains):
    left = Inches(0.75 + i * 4.1)
    box(s, left, Inches(1.7), Inches(3.85), Inches(4.5), WHITE, LINE)
    num = s.shapes.add_textbox(left + Inches(0.3), Inches(2.05), Inches(3.2), Inches(0.35))
    text(num, f"0{i+1}", size=14, bold=True, color=GREEN_DK)
    title = s.shapes.add_textbox(left + Inches(0.3), Inches(2.55), Inches(3.2), Inches(0.7))
    text(title, h, size=20, bold=True, color=INK)
    body = s.shapes.add_textbox(left + Inches(0.3), Inches(3.4), Inches(3.2), Inches(2.2))
    text(body, d, size=15, color=MUTED)
foot(s, 2)
notes(
    s,
    "Say this out loud: we are solving “make headcount recon understandable,” not “collect more HR data.” "
    "Rooted in a real Utah Head of Finance workflow and follow-up review.",
)


# ---------- 3 Why clients ----------
s = slide()
head(
    s,
    "Why this helps clients",
    "When Finance and HR share one headcount story, the business moves faster.",
)
rows = [
    (
        "Finance",
        "Answer “where are we vs the approved hiring plan?” in one view",
        "Fewer recon threads · faster close · board-ready variance",
    ),
    (
        "HR",
        "Log backfill / not backfilling / pivot once — Finance sees it instantly",
        "Less spreadsheet ping-pong · clearer hiring priorities",
    ),
    (
        "Managers",
        "See request status without chasing HR or FP&A",
        "Less tribal knowledge · fewer stalled reqs",
    ),
    (
        "Leadership",
        "Trust the as-of date and the named drivers behind every gap",
        "Decisions on people cost — not debates about the math",
    ),
]
for i, (who, need, out) in enumerate(rows):
    y = Inches(1.6 + i * 1.2)
    box(s, Inches(0.75), y, Inches(2.5), Inches(1.05), SOFT if i % 2 == 0 else WHITE, LINE)
    w = s.shapes.add_textbox(Inches(0.95), y + Inches(0.3), Inches(2.1), Inches(0.45))
    text(w, who, size=16, bold=True, color=INK)
    n = s.shapes.add_textbox(Inches(3.5), y + Inches(0.2), Inches(5.3), Inches(0.7))
    text(n, need, size=15, color=INK)
    o = s.shapes.add_textbox(Inches(8.9), y + Inches(0.2), Inches(3.7), Inches(0.7))
    text(o, out, size=14, bold=True, color=GREEN_DK)
foot(s, 3)
notes(s, "Client why = time, trust, shared language. Tie every prototype feature back to one of these rows.")


# ---------- 4 Why BambooHR ----------
s = slide()
head(
    s,
    "Why this helps BambooHR",
    "A product Finance will fight for — not only HR.",
)
cards = [
    (
        "Sell Finance easier",
        "Head of Finance: “If this product existed, you would sell the finance department much easier.” Direct GTM signal.",
    ),
    (
        "Expand beyond HRIS admin",
        "Headcount recon sits on BambooHR’s source-of-truth. Makes the HRIS the system Finance opens at close — not a spreadsheet export.",
    ),
    (
        "Differentiate on trust",
        "Competitors ship workforce charts. The gap is explained variance + shared dispositions + a clear as-of — hard to copy without HR+Finance design.",
    ),
    (
        "Retention & expansion",
        "Once FP&A depends on BambooHR for the monthly story, switching costs rise and seats expand into finance buyers.",
    ),
]
for i, (h, d) in enumerate(cards):
    col = i % 2
    row = i // 2
    left = Inches(0.75 + col * 6.2)
    top = Inches(1.65 + row * 2.5)
    box(s, left, top, Inches(5.95), Inches(2.25), WHITE, LINE)
    a = s.shapes.add_textbox(left + Inches(0.35), top + Inches(0.35), Inches(5.25), Inches(0.45))
    text(a, h, size=18, bold=True, color=INK)
    b = s.shapes.add_textbox(left + Inches(0.35), top + Inches(0.95), Inches(5.25), Inches(1.0))
    text(b, d, size=14, color=MUTED)
foot(s, 4)
notes(
    s,
    "BambooHR why = GTM into Finance + stickiness of the close ritual. Quote is primary research — treat as directional, not a win-rate claim.",
)


# ---------- 5 What it looks like now ----------
s = slide()
head(
    s,
    "What it looks like today",
    "A multi-tab workbook nobody outside FP&A can navigate.",
    "Fictional Northline data — same structure as a real recon (approved plan · roll-forward · HC recon · net-new).",
)
if SPREADSHEET.exists():
    pic = s.shapes.add_picture(str(SPREADSHEET), Inches(0.7), Inches(1.55), width=Inches(11.9))
    max_h = Inches(5.2)
    if pic.height > max_h:
        pic.width = int(pic.width * (max_h / pic.height))
        pic.height = max_h
else:
    box(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(4.8), SOFT, LINE)
    miss = s.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(10), Inches(0.5))
    text(miss, "Spreadsheet screenshot missing — run deck screenshot generator.", size=16, color=MUTED)
foot(s, 5)
notes(
    s,
    "Emphasize: proprietary contact data was scrubbed. Structure is the point — 4 tabs, dense role rows, monthly columns, formulas. "
    "This is why “presentation tools are complicated and nobody understands them.”",
)


# ---------- 6 Evidence — quotes ----------
s = slide()
head(s, "Evidence", "Primary research — Head of Finance")
box(s, Inches(0.75), Inches(1.7), Inches(12), Inches(2.2), SOFT)
q1 = s.shapes.add_textbox(Inches(1.15), Inches(1.95), Inches(11.2), Inches(1.7))
tf = text(
    q1,
    "“Reconciliation is critical, but the presentation tools are complicated and nobody understands them”",
    size=20,
    color=INK,
)
para(tf, "— Head of Finance", size=14, bold=True, color=GREEN_DK, before=14)

box(s, Inches(0.75), Inches(4.2), Inches(12), Inches(2.2), WHITE, LINE)
q2 = s.shapes.add_textbox(Inches(1.15), Inches(4.45), Inches(11.2), Inches(1.7))
tf = text(
    q2,
    "“If this product existed, you would sell the finance department much easier”",
    size=20,
    color=INK,
)
para(tf, "— Head of Finance  ·  Direct GTM signal for BambooHR", size=14, bold=True, color=GREEN_DK, before=14)
foot(s, 6)
notes(s, "Pain = presentation failure. Opportunity = finance buying motion.")


# ---------- 7 Market — cohesion + recon time ----------
s = slide()
head(
    s,
    "Evidence",
    "Market pattern — HR and Finance still don’t share one headcount system",
)
points = [
    (
        "72%",
        "of HR & Finance teams lack shared systems for workforce planning",
        "Gartner, cited via Kinnect",
    ),
    (
        "14 → 8 days",
        "monthly HR–Finance headcount reconciliation after HRIS–FP&A integration (−43%)",
        "IJIRMPS healthcare case study, 2025",
    ),
    (
        "25%",
        "average discrepancy between HR and Finance headcount reports without a unified structure",
        "PwC (2022), cited in IJIRMPS 2025",
    ),
]
for i, (big, mid, src) in enumerate(points):
    left = Inches(0.75 + i * 4.1)
    box(s, left, Inches(1.7), Inches(3.85), Inches(4.5), WHITE, LINE)
    b = s.shapes.add_textbox(left + Inches(0.25), Inches(2.1), Inches(3.35), Inches(0.9))
    text(b, big, size=36, bold=True, color=GREEN_DK)
    m = s.shapes.add_textbox(left + Inches(0.25), Inches(3.2), Inches(3.35), Inches(1.8))
    text(m, mid, size=15, color=INK)
    c = s.shapes.add_textbox(left + Inches(0.25), Inches(5.3), Inches(3.35), Inches(0.55))
    text(c, src, size=11, color=MUTED)
foot(s, 7)
notes(
    s,
    "Prefer cohesion + recon-time stats. Cite carefully: Gartner via Kinnect; IJIRMPS for 14→8 days and 45% recon-time cut; PwC 25% discrepancy. "
    "Also available: Workday customer claim 60% less manual recon / 50% faster close — use only if you want vendor-sourced color.",
)


# ---------- 8 Hypotheses ----------
s = slide()
head(s, "Hypotheses", "Five beliefs we designed the prototype to test")
hyps = [
    ("H1", "Cognitive load", "Multi-tab recon exceeds working memory", "Story + named gap drivers"),
    ("H2", "Definition debt", "Backfill / not backfilling / pivot undefined", "Shared disposition tags"),
    ("H3", "No decision object", "Mid-cycle changes live only in spreadsheets", "Manual update in-product"),
    ("H4", "Time blindness", "Today vs close date confused", "Explicit July close as-of"),
    ("H5", "Wrong default viz", "Grids serve authors, not leaders", "Bridge tells the story first"),
]
for i, (tag, title, belief, test) in enumerate(hyps):
    y = Inches(1.55 + i * 0.95)
    box(s, Inches(0.75), y, Inches(1.0), Inches(0.75), GREEN)
    tg = s.shapes.add_textbox(Inches(0.75), y + Inches(0.2), Inches(1.0), Inches(0.4))
    text(tg, tag, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    a = s.shapes.add_textbox(Inches(2.0), y + Inches(0.1), Inches(3.5), Inches(0.6))
    text(a, title, size=16, bold=True, color=INK)
    b = s.shapes.add_textbox(Inches(5.6), y + Inches(0.1), Inches(3.8), Inches(0.6))
    text(b, belief, size=14, color=MUTED)
    c = s.shapes.add_textbox(Inches(9.5), y + Inches(0.1), Inches(3.2), Inches(0.6))
    text(c, test, size=14, bold=True, color=GREEN_DK)
foot(s, 8)
notes(s, "AI helped generate; human kept five and mapped each to a UI test.")


# ---------- 9 How I work with the company ----------
s = slide()
head(
    s,
    "How I work with the company",
    "Product owns the problem. Partners own the constraints.",
    "How I’d pull Engineering, Implementation, Design, Finance, and HR into one build — not a handoff chain.",
)
partners = [
    (
        "Product",
        "Frame the problem, cut scope, keep the decision log. Protect the 30-second story.",
    ),
    (
        "Finance (customer + internal FP&A)",
        "Define approved plan vs working plan vs filled. Validate as-of rules and variance language.",
    ),
    (
        "HR / People ops",
        "Own disposition definitions (backfill, not backfilling, pivot) and the approval path.",
    ),
    (
        "Design",
        "Comprehension first. Kill jargon. Test that a new user can read the gap in one pass.",
    ),
    (
        "Engineering",
        "Source of truth, audit trail, and what can be automated vs what must stay a manual update.",
    ),
    (
        "Implementation / CS",
        "How customers adopt at close: glossary, migration from spreadsheets, pilot success criteria.",
    ),
]
for i, (h, d) in enumerate(partners):
    col = i % 3
    row = i // 3
    left = Inches(0.75 + col * 4.1)
    top = Inches(1.65 + row * 2.5)
    box(s, left, top, Inches(3.85), Inches(2.25), WHITE, LINE)
    a = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), Inches(3.35), Inches(0.45))
    text(a, h, size=16, bold=True, color=INK)
    b = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.9), Inches(3.35), Inches(1.1))
    text(b, d, size=14, color=MUTED)
foot(s, 9)
notes(
    s,
    "Speak to collaboration style: weekly triad (Product + Eng + Design), office hours with Finance/HR for glossary, "
    "Implementation in before GA so spreadsheet migration is a product requirement — not a support surprise.",
)


# ---------- 10 Interviews ----------
s = slide()
head(s, "Interviews", "What we have — and what we would learn next")
box(s, Inches(0.75), Inches(1.65), Inches(5.9), Inches(4.7), WHITE, LINE)
t1 = s.shapes.add_textbox(Inches(1.1), Inches(1.95), Inches(5.2), Inches(0.4))
text(t1, "In hand", size=14, bold=True, color=GREEN_DK)
b1 = s.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.2), Inches(3.5))
tf = text(b1, "Head of Finance conversation (exact quotes)", size=15, color=INK)
para(tf, "Live multi-tab workbook: approved plan / roll-forward / HC recon / net-new", size=15, color=INK, before=12)
para(
    tf,
    "Follow-up: name what’s behind the gap, clear as-of, not backfilling + pivots, drop contractors, tell the story, manual mid-cycle edits",
    size=14,
    color=INK,
    before=12,
)

box(s, Inches(6.9), Inches(1.65), Inches(5.7), Inches(4.7), SOFT)
t2 = s.shapes.add_textbox(Inches(7.25), Inches(1.95), Inches(5.1), Inches(0.4))
text(t2, "Next five interviews", size=14, bold=True, color=GREEN_DK)
b2 = s.shapes.add_textbox(Inches(7.25), Inches(2.5), Inches(5.1), Inches(3.5))
tf = text(b2, "FP&A lead — close ritual & recon hours", size=15, color=INK)
para(tf, "HR ops — disposition SLA & ticket path", size=15, color=INK, before=12)
para(tf, "Hiring manager — request status needs", size=15, color=INK, before=12)
para(tf, "Controller / CFO — board report bar", size=15, color=INK, before=12)
para(tf, "Implementation lead — spreadsheet migration pain", size=15, color=INK, before=12)
foot(s, 10)
notes(s, "Honesty: one deep signal + artifact now; structured interviews before scale.")


# ---------- 11 Iterations ----------
s = slide()
head(s, "Iterations", "Ship → show Finance → simplify")
shots = [
    (SHOTS / "01-home-outlook-tile.png", "Story first", "Approved plan · filled · named gap"),
    (SHOTS / "02-bridge-tile.png", "Bridge", "Backfill · not backfilling · pivot"),
    (SHOTS / "03-approvals-tile.png", "Manual update", "Mid-cycle edits systems miss"),
    (SHOTS / "04-export-tile.png", "Export", "Close-ready share-out"),
]
for i, (path, label, sub) in enumerate(shots):
    col = i % 2
    row = i // 2
    left = Inches(0.6 + col * 6.35)
    top = Inches(1.5 + row * 2.7)
    box(s, left, top, Inches(6.1), Inches(0.4), SOFT)
    cap = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.05), Inches(5.7), Inches(0.3))
    text(cap, f"{label}  ·  {sub}", size=12, bold=True, color=INK)
    if path.exists():
        pic = s.shapes.add_picture(str(path), left + Inches(0.05), top + Inches(0.45), width=Inches(6.0))
        max_h = Inches(2.05)
        if pic.height > max_h:
            pic.width = int(pic.width * (max_h / pic.height))
            pic.height = max_h
foot(s, 11)
notes(s, "Partner feedback drove cuts: contractors out; named Marketing −2; pivots; July close vs viewing date; manual update.")


# ---------- 12 Success ----------
s = slide()
head(s, "Success metrics", "Measure understanding — tie every metric to the problem")
metrics = [
    ("Time to understand HC @ close date", "< 30 seconds", "Comprehension"),
    ("Month-end recon time", "Hours → minutes", "Finance speed"),
    ("Disposition logged on attrition", "% of exits tagged", "HR–Finance cohesion"),
    ("Unexplained dept gaps at close", "→ 0", "Named drivers"),
    ("Variance follow-ups / close", "Fewer threads", "Trust"),
    ("Weekly active Finance + HR", "Both personas", "Real use"),
]
for i, (h, target, why) in enumerate(metrics):
    col = i % 3
    row = i // 3
    left = Inches(0.75 + col * 4.1)
    top = Inches(1.65 + row * 2.45)
    box(s, left, top, Inches(3.85), Inches(2.2), WHITE, LINE)
    green = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.85), Inches(0.08))
    green.fill.solid()
    green.fill.fore_color.rgb = GREEN
    green.line.fill.background()
    a = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.35), Inches(3.35), Inches(0.7))
    text(a, h, size=15, bold=True, color=INK)
    b = s.shapes.add_textbox(left + Inches(0.25), top + Inches(1.1), Inches(3.35), Inches(0.4))
    text(b, target, size=18, bold=True, color=GREEN_DK)
    c = s.shapes.add_textbox(left + Inches(0.25), top + Inches(1.6), Inches(3.35), Inches(0.35))
    text(c, why, size=13, color=MUTED)
foot(s, 12)
notes(s, "North star: time to understand headcount at a given close date.")


# ---------- 13 Illustrative impact ----------
s = slide()
head(s, "Illustrative impact", "Modeled for a ~300-person company — labeled assumptions")
note = s.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(12), Inches(0.4))
text(
    note,
    "Not a customer ROI claim — a planning model for interview discussion. Validate hours in the next interviews.",
    size=13,
    bold=True,
    color=MUTED,
)
cards = [
    ("8–12 hrs", "FP&A time / month on recon\nbefore (assumption)", "Matches multi-tab close ritual"),
    ("< 30 min", "Target time to a trusted view\nwith Control Tower", "North-star product goal"),
    ("~40–60%", "Recon-time reductions seen\nin HRIS–FP&A integration cases", "IJIRMPS / industry pattern — validate in pilot"),
]
for i, (big, mid, src) in enumerate(cards):
    left = Inches(0.75 + i * 4.1)
    box(s, left, Inches(2.1), Inches(3.85), Inches(4.0), WHITE, LINE)
    b = s.shapes.add_textbox(left + Inches(0.3), Inches(2.5), Inches(3.25), Inches(0.8))
    text(b, big, size=36, bold=True, color=GREEN_DK)
    m = s.shapes.add_textbox(left + Inches(0.3), Inches(3.5), Inches(3.25), Inches(1.3))
    text(m, mid, size=15, color=INK)
    c = s.shapes.add_textbox(left + Inches(0.3), Inches(5.2), Inches(3.25), Inches(0.6))
    text(c, src, size=12, color=MUTED)
foot(s, 13)
notes(s, "Stay humble on ROI. Point back to market recon-time evidence + partner’s lived spreadsheet pain.")


# ---------- 14 Close ----------
s = slide()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), Inches(7.5))
accent.fill.solid()
accent.fill.fore_color.rgb = GREEN
accent.line.fill.background()
if LOGO.exists():
    s.shapes.add_picture(str(LOGO), Inches(0.9), Inches(1.3), width=Inches(2.2))
t = s.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11), Inches(1.5))
tf = text(t, "Solve the story gap.\nMake headcount recon understandable in 30 seconds.", size=28, bold=True, color=INK)
b = s.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(11), Inches(1.6))
tf = text(b, "Clients get trust and speed. BambooHR gets a Finance buyer who needs the HRIS at close.", size=16, color=INK)
para(tf, "Problem → client why → BambooHR why → evidence → build with partners → measure.", size=16, color=INK, before=10)
para(tf, "AI accelerated discovery. Judgment — and Finance partner review — drove the calls.", size=16, bold=True, color=GREEN_DK, before=10)
c = s.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11), Inches(0.4))
text(c, "Live prototype in README  ·  Next: 5 interviews + one close pilot", size=14, color=MUTED)
notes(s, "End on the problem we solve + who wins. Offer the demo.")

out = ROOT / "Headcount_Control_Tower_Case_Study.pptx"
prs.save(out)
print("Wrote", out, "slides", TOTAL)
